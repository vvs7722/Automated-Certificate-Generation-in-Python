import tkinter as tk
from tkinter import filedialog, messagebox
from pptx import Presentation
import pandas as pd
import os

class CertificateGenerator:
    def __init__(self, root):
        self.root = root
        self.root.title("Certificate Generator")

        self.excel_path = tk.StringVar()
        self.pptx_path = tk.StringVar()
        self.destination_folder = tk.StringVar()

        self.create_widgets()

    def create_widgets(self):
        # Labels to display selected file paths
        excel_label = tk.Label(self.root, textvariable=self.excel_path, wraplength=400)
        excel_label.pack(pady=10)

        pptx_label = tk.Label(self.root, textvariable=self.pptx_path, wraplength=400)
        pptx_label.pack(pady=10)

        # Buttons to trigger file selections
        excel_button = tk.Button(self.root, text="Select Excel File", command=lambda: self.select_file("Excel"))
        excel_button.pack(pady=10)

        pptx_button = tk.Button(self.root, text="Select PPTX File", command=lambda: self.select_file("PPTX"))
        pptx_button.pack(pady=10)

        destination_button = tk.Button(self.root, text="Select Destination Folder", command=self.select_folder)
        destination_button.pack(pady=10)

        merge_button = tk.Button(self.root, text="Start Merging", command=self.merge)
        merge_button.pack(pady=20)

    def select_folder(self):
        folder_path = filedialog.askdirectory(title="Select Folder")
        if folder_path:
            self.destination_folder.set(folder_path)

    def select_file(self, file_type):
        file_path = filedialog.askopenfilename(title=f"Select {file_type} file")
        if file_path:
            if file_type.lower() == 'excel':
                self.excel_path.set(f"Excel file selected: {file_path}")
            elif file_type.lower() == 'pptx':
                self.pptx_path.set(f"PPTX file selected: {file_path}")
            else:
                self.excel_path.set(f"Unknown file type: {file_type}")

    def merge(self):
        # Check if both Excel and PPTX files are selected
        if not self.excel_path.get() or not self.pptx_path.get():
            messagebox.showerror("Error", "Please select both Excel and PPTX files.")
            return

        # Get the string value of the destination_folder StringVar
        destination_folder = self.destination_folder.get()

        # Replace 'your_destination_folder' with the actual destination folder path
        # destination_folder = r'C:\Users\velpo\OneDrive\Desktop\test'

        excel_headings, _ = self.extract_row_names(self.excel_path.get().split(": ")[1])
        print("Excel Headings:", excel_headings)

        df = pd.read_excel(self.excel_path.get().split(": ")[1])

        for index, row in df.iterrows():
            # Extract properties from the first slide
            slide_properties = {}
            for heading in excel_headings:
                slide_properties[heading] = str(row[heading])

            # Create a new presentation with the extracted properties
            modified_presentation = self.extract_and_replace_properties(
                self.pptx_path.get().split(": ")[1], slide_properties
            )

            # Create the destination folder if it doesn't exist
            if not os.path.exists(destination_folder):
                os.makedirs(destination_folder)

            # Save the modified PowerPoint file to the destination folder
            modified_pptx_file = os.path.join(destination_folder, f'certificate_{index+1}.pptx')
            modified_presentation.save(modified_pptx_file)
            print(f"Modified PowerPoint file saved to: {modified_pptx_file}")

        messagebox.showinfo("Success", "Merging completed successfully.")

    @staticmethod
    def extract_row_names(selected_file_path):
        try:
            df = pd.read_excel(selected_file_path)
            sheet_heading = df.columns.values.tolist()
            return sheet_heading, df.shape[0]
        except Exception as e:
            print(f"Error: {e}")
            return [], 0

    @staticmethod
    def extract_and_replace_properties(pptx_path, row_data):
        presentation = Presentation(pptx_path)

        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text_frame = shape.text_frame

                    # Replace old text with new text in the text frame
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            for match_key, match_value in row_data.items():
                                if f'<<{match_key}>>' in run.text:
                                    run.text = run.text.replace(f'<<{match_key}>>', str(match_value))

        return presentation

# Create the main window
root = tk.Tk()
app = CertificateGenerator(root)

# Start the Tkinter event loop
root.mainloop()
