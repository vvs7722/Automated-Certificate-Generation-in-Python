import tkinter as tk
from tkinter import filedialog, messagebox
from email.message import EmailMessage
import smtplib
from pptx import Presentation
import pandas as pd
import os

class CertificateGenerator:
    def __init__(self, root):
        self.root = root
        self.root.title("Certificate Generator")
        self.root.geometry("600x600")
        self.excel_path = tk.StringVar()
        self.pptx_path = tk.StringVar()
        self.destination_folder = tk.StringVar()
        self.sender_email = tk.StringVar()
        self.sender_password = tk.StringVar()
        self.file_name_entry = tk.StringVar()
        self.file_name = ""  # Initialize file_name attribute

        self.create_widgets()

    def create_widgets(self):
        # Title
        title_label = tk.Label(self.root, text="Certificate Generator", font=("Helvetica", 16, "bold"))
        title_label.pack(pady=10)

        # Email Entry
        frame_email = tk.Frame(self.root)
        frame_email.pack(pady=10)
        label_email = tk.Label(frame_email, text="Enter Gmail:", font=("Helvetica", 12, "bold"))
        label_email.grid(row=0, column=0, pady=5)
        entry_email = tk.Entry(frame_email, textvariable=self.sender_email, width=30)
        entry_email.grid(row=0, column=1, pady=5)

        # Password Entry
        frame_password = tk.Frame(self.root)
        frame_password.pack(pady=10)
        label_password = tk.Label(frame_password, text="Enter App Password:", font=("Helvetica", 12, "bold"))
        label_password.grid(row=0, column=0, pady=5)
        entry_password = tk.Entry(frame_password, textvariable=self.sender_password, show="*", width=30)
        entry_password.grid(row=0, column=1, pady=5)

        # File Name Entry
        frame_file_name = tk.Frame(self.root)
        frame_file_name.pack(pady=10)
        label_file_name = tk.Label(frame_file_name, text="Enter Desired File Name:", font=("Helvetica", 12, "bold"))
        label_file_name.grid(row=0, column=0, pady=5)
        entry_file_name = tk.Entry(frame_file_name, textvariable=self.file_name_entry, width=30)
        entry_file_name.grid(row=0, column=1, pady=5)

        # Labels to display selected file paths
        label_excel = tk.Label(self.root, textvariable=self.excel_path, wraplength=400)
        label_excel.pack(pady=5)

        label_pptx = tk.Label(self.root, textvariable=self.pptx_path, wraplength=400)
        label_pptx.pack(pady=5)

        # Buttons to trigger file selections
        button_excel = tk.Button(self.root, text="Select Excel File", command=lambda: self.select_file("Excel"))
        button_excel.pack(pady=10)

        button_pptx = tk.Button(self.root, text="Select PPTX File", command=lambda: self.select_file("PPTX"))
        button_pptx.pack(pady=10)

        button_destination = tk.Button(self.root, text="Select Destination Folder", command=self.select_folder)
        button_destination.pack(pady=10)

            # Confirm Button
        button_confirm = tk.Button(self.root, text="Confirm", command=self.get_entry_text, bg="#4CAF50", fg="white", font=("Helvetica", 12, "bold"))
        button_confirm.pack(pady=10)

        # Start Merging Button
        button_merge = tk.Button(self.root, text="Start Merging", command=self.merge, bg="#008CBA", fg="white", font=("Helvetica", 12, "bold"))
        button_merge.pack(pady=20)
    def select_folder(self):
        folder_path = filedialog.askdirectory(title="Select Destination Folder")
        if folder_path:
            self.destination_folder.set(folder_path)

    def select_file(self, file_type):
        file_path = filedialog.askopenfilename(title=f"Select {file_type} File")
        if file_path:
            if file_type.lower() == 'excel':
                self.excel_path.set(f"Excel File Selected: {file_path}")
            elif file_type.lower() == 'pptx':
                self.pptx_path.set(f"PPTX File Selected: {file_path}")
            else:
                self.excel_path.set(f"Unknown File Type: {file_type}")

    def get_entry_text(self):
        self.sender_email_value = self.sender_email.get()
        self.sender_password_value = self.sender_password.get()
        file_name = self.file_name_entry.get()
        self.file_name = file_name.replace(" ", "-")


    def merge(self):
        # Check if both Excel and PPTX files are selected
        if not self.excel_path.get() or not self.pptx_path.get():
            messagebox.showerror("Error", "Please select both Excel and PPTX files.")
            return

        # Get the string value of the destination_folder StringVar
        destination_folder = self.destination_folder.get()

        excel_headings, _ = self.extract_row_names(self.excel_path.get().split(": ")[1])
        print("Excel Headings:", excel_headings)

        df = pd.read_excel(self.excel_path.get().split(": ")[1])
        email_addresses = df["email"].tolist()

        for index, row in df.iterrows():
            # Extract properties from the first slide
            slide_properties = {}
            for heading in excel_headings:
                slide_properties[heading] = str(row[heading])

            # Create a new presentation with the extracted properties
            modified_presentation = self.extract_and_replace_properties(
                self.pptx_path.get().split(": ")[1], slide_properties
            )

            # Create the destination folder if it doesn't exist
            if not os.path.exists(destination_folder):
                os.makedirs(destination_folder)

            # Save the modified PowerPoint file to the destination folder
            modified_pptx_file = os.path.join(destination_folder, f'{self.file_name}_certificate.pptx')
            modified_presentation.save(modified_pptx_file)

            # Send email with the modified PowerPoint file as an attachment
            attachment_path = modified_pptx_file
            self.send_email(attachment_path, email_addresses[index])

            print(f"Modified PowerPoint file saved to: {modified_pptx_file}")

        messagebox.showinfo("Success", "Merging completed successfully.")

    def extract_row_names(self, selected_file_path):
        try:
            df = pd.read_excel(selected_file_path)
            sheet_heading = df.columns.values.tolist()
            return sheet_heading, df.shape[0]
        except Exception as e:
            print(f"Error: {e}")
            return [], 0

    def extract_and_replace_properties(self, pptx_path, row_data):
        presentation = Presentation(pptx_path)

        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text_frame = shape.text_frame

                    # Replace old text with new text in the text frame
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            for match_key, match_value in row_data.items():
                                if f'<<{match_key}>>' in run.text:
                                    run.text = run.text.replace(f'<<{match_key}>>', str(match_value))

        return presentation

    def send_email(self, attachment_path, receiver_email):
        subject = 'Presentation Attached'
        body = 'Please find the attached PowerPoint presentation.'
        message = EmailMessage()
        message.set_content(body)
        message['Subject'] = subject
        message['From'] = self.sender_email_value
        message['To'] = receiver_email

        # Attach the PowerPoint file
        with open(attachment_path, 'rb') as attachment:
            message.add_attachment(attachment.read(), maintype='application',
                                   subtype='vnd.openxmlformats-officedocument.presentationml.presentation',
                                   filename=f'{self.file_name}_certificate.pptx')

        # Connect to the SMTP server and send the email
        try:
            server = smtplib.SMTP('smtp.gmail.com', 587)
            server.starttls()
            server.login(self.sender_email_value, self.sender_password_value)
            server.send_message(message)
            server.quit()
            print("Email sent successfully!")
        except Exception as e:
            print(f"Error: {e}")

# Create the main window
root = tk.Tk()
app = CertificateGenerator(root)

# Start the Tkinter event loop
root.mainloop()
